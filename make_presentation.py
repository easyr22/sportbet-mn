"""SportBet MN — Танилцуулга PowerPoint үүсгэгч"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dgm.color import RGBColor as _RGB  # noqa: F401  # placeholder
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# 1xBet-style colors
BG_DARK     = RGBColor(0x0E, 0x1B, 0x2C)
BG_PANEL    = RGBColor(0x1A, 0x2C, 0x42)
BLUE        = RGBColor(0x1E, 0x5A, 0x99)
BLUE_DARK   = RGBColor(0x15, 0x49, 0x7F)
GREEN       = RGBColor(0x22, 0xB1, 0x4C)
ORANGE      = RGBColor(0xFF, 0xA5, 0x00)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_GREY   = RGBColor(0xA8, 0xB5, 0xC7)
YELLOW      = RGBColor(0xFF, 0xCC, 0x00)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

def fill_bg(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    return bg

def add_text(slide, x, y, w, h, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left=tf.margin_right=tf.margin_top=tf.margin_bottom=Pt(2)
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = font
    return tb

def add_box(slide, x, y, w, h, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None: s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(1)
    s.shadow.inherit = False
    return s

def add_accent_bar(slide, x, y, h, color=GREEN, w=Inches(0.08)):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background(); s.shadow.inherit = False

# ───────────── SLIDE 1: TITLE ─────────────
s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
fill_bg(s, BG_DARK)
# Top gradient panel
top = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(2.2))
top.fill.solid(); top.fill.fore_color.rgb = BLUE_DARK
top.line.fill.background(); top.shadow.inherit = False

add_text(s, Inches(0.7), Inches(2.7), Inches(12), Inches(1.2),
         "SportBet MN", size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.7), Inches(3.9), Inches(12), Inches(0.6),
         "Спорт бооцооны вэб платформ", size=24, color=TEXT_GREY, align=PP_ALIGN.CENTER)

# MN badge
b = add_box(s, Inches(8.5), Inches(2.85), Inches(0.9), Inches(0.55), GREEN)
add_text(s, Inches(8.5), Inches(2.93), Inches(0.9), Inches(0.5),
         "MN", size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# URL pill
url_box = add_box(s, Inches(3.5), Inches(5.0), Inches(6.3), Inches(0.7), BG_PANEL, BLUE)
add_text(s, Inches(3.5), Inches(5.13), Inches(6.3), Inches(0.5),
         "🌐  https://sportbet-mn.onrender.com", size=18, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s, Inches(0.7), Inches(6.5), Inches(12), Inches(0.5),
         "Бодит цагийн тоглолт  •  Бодит мөнгө бус демо  •  Мобил дэмжлэгтэй",
         size=14, color=ORANGE, align=PP_ALIGN.CENTER)

# ───────────── SLIDE 2: ТӨСЛИЙН ТАНИЛЦУУЛГА ─────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
fill_bg(s, BG_DARK)
add_accent_bar(s, Inches(0.5), Inches(0.5), Inches(0.6))
add_text(s, Inches(0.75), Inches(0.45), Inches(12), Inches(0.7),
         "ТӨСЛИЙН ТАНИЛЦУУЛГА", size=28, bold=True, color=WHITE)

content = [
    ("🎯 Зорилго", "Мэргэжлийн төвшний спорт бооцооны вэб сайт хийх (1xBet-ийн загвараар)"),
    ("⚡ Технологи", "Frontend: Flutter Web  •  Backend: Python (Flask)  •  Hosting: Render.com"),
    ("📊 Өгөгдөл", "ESPN public API + бодит баг/тоглогчдын мэдээлэл (276+ тоглолт, бүгд live)"),
    ("👥 Зорилтот", "Сургуулийн төгсөлтийн төсөл — багшид болон ангийнхандаа танилцуулга хийх"),
    ("🚀 Үр дүн", "Олон улсын домейн, мобил дэмжлэг, имэйл бүртгэл, бооцооны систем"),
]
y = Inches(1.6)
for title, desc in content:
    add_box(s, Inches(0.7), y, Inches(11.9), Inches(0.95), BG_PANEL)
    add_accent_bar(s, Inches(0.7), y, Inches(0.95), color=BLUE)
    add_text(s, Inches(0.95), y + Inches(0.12), Inches(3.2), Inches(0.4),
             title, size=15, bold=True, color=YELLOW)
    add_text(s, Inches(0.95), y + Inches(0.45), Inches(11), Inches(0.5),
             desc, size=13, color=WHITE)
    y += Inches(1.05)

# ───────────── SLIDE 3: ОНЦЛОГ БОЛОМЖУУД ─────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
fill_bg(s, BG_DARK)
add_accent_bar(s, Inches(0.5), Inches(0.5), Inches(0.6))
add_text(s, Inches(0.75), Inches(0.45), Inches(12), Inches(0.7),
         "ОНЦЛОГ БОЛОМЖУУД", size=28, bold=True, color=WHITE)

features = [
    ("🔴", "Бодит цагийн тоглолт", "276+ live тоглолт, минут тутамд шинэчлэгддэг"),
    ("⚽", "11 төрлийн спорт", "Хөл бөмбөг, Сагсан, Эспорт, Теннис, MMA, Хоккей, Бейсбол..."),
    ("🎮", "Эспорт", "CS2, LoL, Dota2, Valorant, Rocket League — бодит баг, тэмцээн"),
    ("📧", "Имэйл бүртгэл", "6 оронтой кодоор баталгаажуулсан аюулгүй бүртгэл"),
    ("💰", "Бооцооны систем", "1X2, Over/Under, Both Teams To Score, Аккумулятор..."),
    ("📱", "PWA дэмжлэг", "iOS/Android утсанд аппликейшн шиг суулгах боломжтой"),
]
cols, rows = 3, 2
cw = Inches(4.0); ch = Inches(2.6); gap = Inches(0.15)
sx = Inches(0.7); sy = Inches(1.5)
for i, (emo, title, desc) in enumerate(features):
    r, c = i // cols, i % cols
    x = sx + (cw + gap) * c
    y = sy + (ch + gap) * r
    add_box(s, x, y, cw, ch, BG_PANEL)
    add_accent_bar(s, x, y, ch, color=GREEN if i % 2 == 0 else ORANGE)
    add_text(s, x + Inches(0.25), y + Inches(0.2), Inches(1), Inches(0.7),
             emo, size=36, color=WHITE)
    add_text(s, x + Inches(0.25), y + Inches(1.0), cw - Inches(0.5), Inches(0.45),
             title, size=15, bold=True, color=YELLOW)
    add_text(s, x + Inches(0.25), y + Inches(1.5), cw - Inches(0.5), Inches(1.0),
             desc, size=11, color=TEXT_GREY)

# ───────────── SLIDE 4: ТЕХНИК ─────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
fill_bg(s, BG_DARK)
add_accent_bar(s, Inches(0.5), Inches(0.5), Inches(0.6))
add_text(s, Inches(0.75), Inches(0.45), Inches(12), Inches(0.7),
         "ТЕХНИК ХЭРЭГЛЭСЭН", size=28, bold=True, color=WHITE)

# Frontend column
add_text(s, Inches(0.7), Inches(1.5), Inches(6), Inches(0.5),
         "FRONTEND", size=18, bold=True, color=GREEN)
fe_box = add_box(s, Inches(0.7), Inches(2.0), Inches(6.0), Inches(4.5), BG_PANEL)
fe_items = [
    "▸  Flutter Web 3.x",
    "▸  Dart програмчлалын хэл",
    "▸  Provider (state management)",
    "▸  Material Design компонент",
    "▸  HTTP client API дуудалт",
    "▸  PWA (Manifest + Service Worker)",
    "▸  Responsive UI (мобил/таблет/десктоп)",
]
y = Inches(2.2)
for it in fe_items:
    add_text(s, Inches(1.0), y, Inches(5.5), Inches(0.45),
             it, size=14, color=WHITE)
    y += Inches(0.55)

# Backend column
add_text(s, Inches(7.0), Inches(1.5), Inches(6), Inches(0.5),
         "BACKEND & HOSTING", size=18, bold=True, color=ORANGE)
be_box = add_box(s, Inches(7.0), Inches(2.0), Inches(5.7), Inches(4.5), BG_PANEL)
be_items = [
    "▸  Python 3.11 + Flask framework",
    "▸  Gunicorn production server",
    "▸  ESPN Sports API (бодит мэдээлэл)",
    "▸  Gmail SMTP (имэйл илгээлт)",
    "▸  Docker контейнер",
    "▸  Render.com cloud hosting",
    "▸  GitHub автомат deploy",
]
y = Inches(2.2)
for it in be_items:
    add_text(s, Inches(7.3), y, Inches(5.3), Inches(0.45),
             it, size=14, color=WHITE)
    y += Inches(0.55)

# ───────────── SLIDE 5: МЭДЭЭЛЛИЙН ЭХ ҮҮСВЭР ─────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
fill_bg(s, BG_DARK)
add_accent_bar(s, Inches(0.5), Inches(0.5), Inches(0.6))
add_text(s, Inches(0.75), Inches(0.45), Inches(12), Inches(0.7),
         "ХАМРАГДСАН ТЭМЦЭЭН", size=28, bold=True, color=WHITE)

sports = [
    ("⚽ Хөл бөмбөг", "30 лиг", "PL, La Liga, Serie A, Bundesliga, Ligue 1, Champions League..."),
    ("🏀 Сагсан бөмбөг", "NBA, WNBA, NCAA", "Boston Celtics, LA Lakers, Warriors, Bucks..."),
    ("⚾ Бейсбол", "MLB, NCAA", "Yankees, Red Sox, Dodgers, Giants..."),
    ("🏒 Хоккей", "NHL Playoffs", "Maple Leafs, Bruins, Oilers, Lightning..."),
    ("🏈 Америк хөл бөмбөг", "NFL, NCAA", "Chiefs, Bills, Eagles, 49ers..."),
    ("🎮 Эспорт", "16 тэмцээн", "BLAST, IEM, PGL Major, LCK, MSI, VCT..."),
    ("🎾 Теннис", "ATP, WTA", "Alcaraz, Sinner, Djokovic, Swiatek, Sabalenka..."),
    ("🥊 MMA", "UFC, PFL", "Makhachev, Pereira, Jones, Edwards..."),
    ("🏐 Бусад", "Volleyball, Бадминтон, Ширээний теннис", "FIVB, BWF, WTT тэмцээнүүд"),
]
y = Inches(1.5)
for title, count, desc in sports:
    add_box(s, Inches(0.7), y, Inches(11.9), Inches(0.55), BG_PANEL)
    add_text(s, Inches(0.95), y + Inches(0.1), Inches(4), Inches(0.4),
             title, size=14, bold=True, color=YELLOW)
    add_text(s, Inches(4.5), y + Inches(0.1), Inches(2.5), Inches(0.4),
             count, size=13, bold=True, color=GREEN)
    add_text(s, Inches(7.0), y + Inches(0.1), Inches(5.7), Inches(0.4),
             desc, size=11, color=TEXT_GREY)
    y += Inches(0.6)

# ───────────── SLIDE 6: АЯЛАЛ ХӨГЖҮҮЛЭЛТ ─────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
fill_bg(s, BG_DARK)
add_accent_bar(s, Inches(0.5), Inches(0.5), Inches(0.6))
add_text(s, Inches(0.75), Inches(0.45), Inches(12), Inches(0.7),
         "ХЭРЭГЛЭГЧИЙН АЯЛАЛ", size=28, bold=True, color=WHITE)

steps = [
    ("1", "Site нээх", "Хэрэглэгч https://sportbet-mn.onrender.com нээгээд тоглолтуудыг харна"),
    ("2", "Бүртгүүлэх", "Имэйл хаяг + хэрэглэгчийн нэр + нууц үг → 6 оронтой код имэйл рүү"),
    ("3", "Баталгаажуулах", "Имэйлээс кодоо хараад оруулна → 50,000₮ урамшуулал автоматаар"),
    ("4", "Тоглолт сонгох", "Бодит live тоглолтоос odds сонгох (1X2, Over/Under, BTS...)"),
    ("5", "Бооцоо тавих", "Bet slip-д цуглуулаад дүн оруулж 'Бооцоо тавих' дарна"),
    ("6", "Үлдэгдэл харах", "Профайл хуудаснаас Deposit/Withdraw хийж, гүйлгээний түүх харна"),
]
y = Inches(1.5)
for num, title, desc in steps:
    # Number circle
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), y, Inches(0.7), Inches(0.7))
    c.fill.solid(); c.fill.fore_color.rgb = GREEN
    c.line.fill.background(); c.shadow.inherit = False
    add_text(s, Inches(0.8), y + Inches(0.13), Inches(0.7), Inches(0.5),
             num, size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Card
    add_box(s, Inches(1.7), y, Inches(10.9), Inches(0.7), BG_PANEL)
    add_text(s, Inches(1.95), y + Inches(0.07), Inches(3), Inches(0.35),
             title, size=14, bold=True, color=YELLOW)
    add_text(s, Inches(1.95), y + Inches(0.36), Inches(10.5), Inches(0.35),
             desc, size=11, color=WHITE)
    y += Inches(0.85)

# ───────────── SLIDE 7: ДЭМО ─────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
fill_bg(s, BG_DARK)
add_accent_bar(s, Inches(0.5), Inches(0.5), Inches(0.6))
add_text(s, Inches(0.75), Inches(0.45), Inches(12), Inches(0.7),
         "ШУУД ДЭМО", size=28, bold=True, color=WHITE)

# QR placeholder + URL
add_box(s, Inches(2.3), Inches(1.7), Inches(9), Inches(4.3), BLUE)
add_text(s, Inches(2.3), Inches(2.0), Inches(9), Inches(0.6),
         "Утсаараа эсвэл компьютероор нээ:", size=18, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(2.3), Inches(2.8), Inches(9), Inches(0.8),
         "https://sportbet-mn.onrender.com", size=32, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)
add_text(s, Inches(2.3), Inches(3.9), Inches(9), Inches(0.5),
         "📱 iPhone Safari → Share → Add to Home Screen",
         size=14, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(2.3), Inches(4.4), Inches(9), Inches(0.5),
         "🤖 Android Chrome → ⋮ цэс → Install app",
         size=14, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(2.3), Inches(5.1), Inches(9), Inches(0.5),
         "⏱ Эхний удаа нээхэд ~30 секунд хүлээх (Render free сервер)",
         size=12, color=ORANGE, align=PP_ALIGN.CENTER)

# ───────────── SLIDE 8: ЦААШДЫН ХӨГЖҮҮЛЭЛТ ─────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
fill_bg(s, BG_DARK)
add_accent_bar(s, Inches(0.5), Inches(0.5), Inches(0.6))
add_text(s, Inches(0.75), Inches(0.45), Inches(12), Inches(0.7),
         "ЦААШДЫН ТӨЛӨВЛӨГӨӨ", size=28, bold=True, color=WHITE)

future = [
    ("💳 QPay интеграц", "Жинхэнэ Монгол төлбөрийн систем — мерчант гэрээтэй холбох"),
    ("🗄️  PostgreSQL ӨГЁС", "Хэрэглэгч, гүйлгээ, бооцоог байнгын мэдээллийн сан руу"),
    ("📊 Live статистик", "Тоглолт бүрд гүн статистик, h2h түүх, тоглогчдын мэдээлэл"),
    ("🎰 Казино, TOTO", "Slot тоглоом, lotto, виртуал спорт"),
    ("🤖 ML odds", "Машин сургалтаар тоглолтын өгөөжийн загвар сайжруулах"),
    ("🌐 Олон хэл", "Англи, Орос хэл нэмэх"),
    ("📱 Native app", "Flutter-ээс iOS/Android жинхэнэ аппликейшн build хийх"),
]
y = Inches(1.5)
for title, desc in future:
    add_box(s, Inches(0.7), y, Inches(11.9), Inches(0.7), BG_PANEL)
    add_accent_bar(s, Inches(0.7), y, Inches(0.7), color=ORANGE)
    add_text(s, Inches(0.95), y + Inches(0.1), Inches(3.5), Inches(0.4),
             title, size=14, bold=True, color=YELLOW)
    add_text(s, Inches(4.5), y + Inches(0.13), Inches(8), Inches(0.4),
             desc, size=12, color=WHITE)
    y += Inches(0.78)

# ───────────── SLIDE 9: БАЯРЛАЛАА ─────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
fill_bg(s, BG_DARK)
# gradient strip
strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.5), prs.slide_width, Inches(2.5))
strip.fill.solid(); strip.fill.fore_color.rgb = BLUE
strip.line.fill.background(); strip.shadow.inherit = False

add_text(s, Inches(0.5), Inches(2.9), Inches(12.3), Inches(1.0),
         "БАЯРЛАЛАА", size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(4.0), Inches(12.3), Inches(0.6),
         "Танилцуулсанд баярлалаа!", size=20, color=YELLOW, align=PP_ALIGN.CENTER)

add_text(s, Inches(0.5), Inches(5.4), Inches(12.3), Inches(0.5),
         "🌐 https://sportbet-mn.onrender.com",
         size=18, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.5),
         "📂 GitHub: easyr22/sportbet-mn",
         size=14, color=TEXT_GREY, align=PP_ALIGN.CENTER)

# Save
out = "C:\\Users\\hitech\\OneDrive\\Desktop\\sportbet_mn2\\SportBet_MN_Tanitsuulga.pptx"
prs.save(out)
print(f"Done: {out}")
